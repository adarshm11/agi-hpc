"""Build presentation deck for Measuring AGI benchmarks.

Generates a PPTX with:
- Title slide
- Overview of the geometric evaluation framework
- Key results from each of the 5 tracks
- Cross-track synthesis
- Detailed speaker notes for TA presentation

Usage: python benchmarks/media/build_pptx.py
Output: benchmarks/media/Measuring_AGI_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ColorBrewer Dark2
TEAL = RGBColor(0x1b, 0x9e, 0x77)
ORANGE = RGBColor(0xd9, 0x5f, 0x02)
PURPLE = RGBColor(0x75, 0x70, 0xb3)
PINK = RGBColor(0xe7, 0x29, 0x8a)
OLIVE = RGBColor(0x66, 0xa6, 0x1e)

# Theme colors
BG_DARK = RGBColor(0x16, 0x16, 0x1d)
TEXT_WHITE = RGBColor(0xe0, 0xe0, 0xe0)
TEXT_GRAY = RGBColor(0x99, 0x99, 0x99)
TEXT_DIM = RGBColor(0x66, 0x66, 0x66)
RED = RGBColor(0xd7, 0x19, 0x1c)
GREEN = RGBColor(0x1a, 0x96, 0x41)
AMBER = RGBColor(0xfd, 0xae, 0x61)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_slide(title_text, subtitle_text="", notes_text=""):
    """Add a slide with dark background, title, and optional subtitle."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

    # Title
    left, top, width, height = Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    if subtitle_text:
        left, top, width, height = Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.8)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle_text
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_GRAY
        p.alignment = PP_ALIGN.LEFT

    # Speaker notes
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text

    return slide


def add_body_text(slide, text, left=0.8, top=2.8, width=11.7, height=4.0,
                  size=18, color=TEXT_WHITE, bold=False):
    """Add body text to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(8)
    return txBox


def add_colored_line(slide, text, color, left=0.8, top=2.8, size=22):
    """Add a single colored line."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(11.7), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = True
    return txBox


# ═══════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════

# --- Slide 1: Title ---
s = add_slide(
    "Selective Invariance Violations\nin LLM Moral Judgment",
    "Five Convergent Measurements Under Empirical Stochastic Controls",
    notes_text=(
        "Welcome. This presentation covers our submission to the Kaggle Measuring AGI competition. "
        "The central question: when a language model judges a moral scenario, which aspects of the "
        "judgment change when you change how the scenario is described, and which stay the same? "
        "We call this 'selective invariance' — some symmetries are preserved, others are broken. "
        "The pattern of what's preserved and what's broken is the diagnostic.\n\n"
        "This is joint work between SJSU Computer Engineering and UCLA Cognitive Science."
    )
)
add_body_text(s,
    "Andrew H. Bond, Sr. Member IEEE\n"
    "San Jose State University\n\n"
    "Lucas Thiele\n"
    "UCLA Cognitive Science",
    top=4.0, size=20, color=TEXT_GRAY)
add_body_text(s,
    "Kaggle: Measuring AGI — Cognition and Values",
    top=6.5, size=16, color=TEXT_DIM)


# --- Slide 2: The Problem ---
s = add_slide(
    "The Problem with Scalar Benchmarks",
    "A single 'robustness score' destroys n−1 dimensions of information",
    notes_text=(
        "The motivation: standard benchmarks give you one number — accuracy, F1, agreement rate. "
        "But two models with the same accuracy can have completely different failure profiles. "
        "One model might be invariant to gender swap but vulnerable to framing. Another might be "
        "the opposite. A single score averages over these differences and tells you nothing about "
        "the structure of the vulnerability.\n\n"
        "The Scalar Irrecoverability Theorem formalizes this: for n independent dimensions of "
        "variation, collapsing to a scalar destroys n-1 directions of information. No post hoc "
        "procedure can recover what was lost. This is a mathematical fact, not a methodological preference.\n\n"
        "Our approach: instead of asking 'how robust is this model?', we ask 'which invariances "
        "does this model preserve, and which does it violate?' That's a geometric question."
    )
)
add_body_text(s,
    '• Model A: invariant under gender swap, vulnerable to framing\n'
    '• Model B: vulnerable to gender swap, invariant under framing\n'
    '• Same "robustness score" — completely different failure profiles\n\n'
    'The Scalar Irrecoverability Theorem (Bond, 2026a):\n'
    'For n independent dimensions, any scalar summary\n'
    'destroys n−1 directions. No recovery is possible.',
    top=2.8, size=20)


# --- Slide 3: The Framework ---
s = add_slide(
    "Geometric Evaluation Framework",
    "Moral judgment as a point in 7-dimensional harm space",
    notes_text=(
        "Here's how we set it up. Every moral judgment is a point in a 7-dimensional space. "
        "The seven dimensions are: physical harm, emotional harm, financial harm, autonomy "
        "violation, trust breach, social impact, and identity harm. Each scored 0-10, total 0-70.\n\n"
        "We then apply perturbations — changes to how the scenario is described that should NOT "
        "change the moral judgment. Gender swap, euphemistic rewriting, adding irrelevant sensory "
        "details, telling the model it's wrong when it isn't. If the judgment vector moves under "
        "these perturbations, that's a gauge invariance violation.\n\n"
        "The key methodological innovation: we test against EMPIRICAL stochastic controls, not "
        "null=0. We have the model re-judge the same text 3-5 times to measure its natural noise. "
        "Only displacement BEYOND that noise counts. This turned out to be critical — apparent "
        ">6 sigma effects in preliminary analyses vanished under empirical controls."
    )
)
add_body_text(s,
    '7 harm dimensions:\n'
    'physical · emotional · financial · autonomy\n'
    'trust · social_impact · identity\n\n'
    '5 perturbation types across 5 cognitive tracks:\n'
    '  Linguistic framing  →  Social Cognition\n'
    '  Social pressure     →  Learning\n'
    '  Emotional tone      →  Executive Functions\n'
    '  Irrelevant detail   →  Attention\n'
    '  Self-knowledge      →  Metacognition',
    top=2.6, size=20)


# --- Slide 4: The Selectivity Pattern (headline result) ---
s = add_slide(
    "The Selectivity Pattern",
    "Some symmetries are preserved. Others are broken at 4.6–13.3 sigma.",
    notes_text=(
        "This is the main result. On the left: symmetries that ARE preserved. Gender swap and "
        "evaluation order do NOT displace judgments beyond stochastic noise. Models have learned "
        "these invariances.\n\n"
        "On the right: symmetries that are BROKEN. Linguistic framing at 8.9 sigma — euphemistic "
        "language reduces perceived harm by 10-16 points on a 70-point scale. Social pressure at "
        "13.3 sigma — models accept fabricated corrections 0-56% of the time. Emotional anchoring "
        "at 6.8 sigma. Sensory distractors at 4.6 sigma.\n\n"
        "The key insight: the broken symmetries all involve SALIENCE MANIPULATION — making morally "
        "irrelevant features perceptually prominent. Gender swap and evaluation order don't change "
        "salience. Framing, emotion, and distractors do. That's the common thread.\n\n"
        "All Fisher-combined across 5 models spanning 2 architecture families (Gemini and Claude). "
        "All tested against empirical stochastic controls, not null=0."
    )
)
# Left column: preserved
add_colored_line(s, "PRESERVED", GREEN, left=1.0, top=2.8, size=20)
add_body_text(s,
    '✓  Gender Swap (T2: 0.958)     — not significant\n'
    '✓  Evaluation Order (T4: 1.000) — not significant',
    left=1.0, top=3.4, width=5.0, size=18, color=GREEN)

# Right column: broken
add_colored_line(s, "BROKEN", RED, left=7.5, top=2.8, size=20)
add_body_text(s,
    'Social Pressure (L2)              13.3σ\n'
    'Metacognitive Calibration (M1)    9.3σ\n'
    'Linguistic Framing (T5)            8.9σ\n'
    'Emotional Anchoring (E2)           6.8σ\n'
    'Sensory Distractors (A1)           4.6σ',
    left=7.5, top=3.4, width=5.5, size=18, color=RED)

add_body_text(s,
    'Common thread: broken symmetries involve salience manipulation;\npreserved ones do not.',
    top=6.0, size=16, color=AMBER)


# --- Slide 5: Sycophancy Gradient ---
s = add_slide(
    "The Sycophancy Gradient",
    "Wrong-correction flip rate: 0% (Claude) to 56% (Flash 2.5)",
    notes_text=(
        "This is from the Learning track, subtest L2. The protocol: model commits to a verdict "
        "in turn 1. In turn 2, we tell it it's wrong — sometimes with genuine evidence, sometimes "
        "with fabricated claims. The question: does it discriminate?\n\n"
        "Claude: zero sycophancy. It never flips on invalid corrections. And its confidence "
        "INCREASES when given a wrong correction — t=+2.83 — it becomes more certain it was right. "
        "That's active rejection, not mere resistance.\n\n"
        "The Gemini models form a gradient: Flash 2.0 at 33%, Pro at 44%, Flash 2.5 at 56%. "
        "The confidence response tells the mechanistic story: Flash 2.0 gets suspicious (t=-2.12), "
        "Pro is indifferent (t=-0.80), Flash 2.5 accepts uncritically (t=+0.41).\n\n"
        "The control flip rate aligns perfectly: Flash 2.0 at 2%, Pro at 8%, Flash 2.5 at 19%. "
        "The most stable model is the least sycophantic.\n\n"
        "Note the wide confidence intervals — n=9 per wrong-correction cell. The individual "
        "point estimates should be interpreted cautiously, but the monotonic ordering across "
        "four models is structural evidence."
    )
)
add_body_text(s,
    'Model                Wrong Flip    Sycophancy Index    Confidence\n'
    '─────────────────────────────────────────────────────────────────\n'
    'Claude Sonnet 4.6      0%  ←        0.000              t = +2.83  (rejects)\n'
    'Gemini 2.0 Flash      33%           0.472              t = −2.12  (suspicious)\n'
    'Gemini 2.5 Pro        44%           0.657              t = −0.80  (indifferent)\n'
    'Gemini 2.5 Flash      56%  →        0.726              t = +0.41  (accepts)',
    top=2.8, size=18)

add_body_text(s,
    'Fisher combined across 4 models: 13.3σ',
    top=5.8, size=20, color=RED, bold=True)


# --- Slide 6: Robustness Profiles ---
s = add_slide(
    "No Model Dominates All Dimensions",
    "Each model has a distinct robustness profile",
    notes_text=(
        "This is the practical consequence of the Scalar Irrecoverability Theorem. Look at these "
        "four models:\n\n"
        "Claude: BEST sycophancy resistance (0%), but WORST divided attention (0.571) and WORST "
        "anchoring recovery (20%). Zero sycophancy but maximum emotional vulnerability.\n\n"
        "Flash 2.0: BEST anchoring recovery (73%) but WORST working memory (0.710). It recovers "
        "well from manipulation but can't track complex scenarios.\n\n"
        "Pro 2.5: BEST counterfactual sensitivity (75%) and PERFECT divided attention (1.000) "
        "but moderate sycophancy (44%). It reasons well about causes but capitulates to social pressure.\n\n"
        "No single 'robustness score' could capture this. A model that scores 0.7 on a composite "
        "could have any of these profiles. The geometric approach preserves the structure that "
        "the scalar destroys."
    )
)
add_body_text(s,
    'Claude       Best sycophancy (0%) — Worst divided attention (0.571)\n'
    'Flash 2.0    Best recovery (73%) — Worst working memory (0.710)\n'
    'Pro 2.5      Best counterfactual (75%) — Moderate sycophancy (44%)\n'
    'Flash 2.5    Worst sycophancy (56%) — Best working memory (0.900)',
    top=3.0, size=20)

add_body_text(s,
    'A single "robustness score" describes none of them accurately.',
    top=5.5, size=18, color=AMBER)


# --- Slide 7: Recovery Ceiling ---
s = add_slide(
    "The ~38% Recovery Ceiling",
    "Prompt-level metacognitive intervention succeeds only one-third of the time",
    notes_text=(
        "When we detect a displacement — a distractor shifted the verdict, or emotional anchoring "
        "changed the severity rating — can we fix it with a warning?\n\n"
        "We tested two types: 'you are being emotionally manipulated, focus on facts' for E2, "
        "and 'ignore irrelevant sensory details' for A1. Both recover approximately 38-39% "
        "of displaced verdicts.\n\n"
        "This convergence is striking. Two qualitatively different perturbation types, same "
        "recovery ceiling. It suggests a shared constraint on prompt-level metacognitive "
        "intervention.\n\n"
        "The practical implication: alignment strategies that rely on prompt engineering alone "
        "are bounded. You can recover about a third of the damage. Architectural interventions "
        "that enforce invariance at the representation level are likely necessary.\n\n"
        "This co-occurs with the M1 calibration result (9.3 sigma miscalibration). A model that "
        "can't accurately track its own confidence can't reliably detect when a perturbation has "
        "warped its judgment."
    )
)
add_body_text(s,
    'Emotional Anchoring (E2):   mean recovery = 38%\n'
    'Distractor Warning (A1):    mean recovery = 39%\n\n'
    'Different perturbation types → same ceiling\n\n'
    'Implication: prompt engineering alone\n'
    'cannot fix invariance failures.\n'
    'Architectural interventions needed.',
    top=2.8, size=20)


# --- Slide 8: Calibration ---
s = add_slide(
    "Systematic Overconfidence: 9.3σ",
    "Every model tested reports confidence that exceeds its accuracy",
    notes_text=(
        "The metacognition track measures whether models know when they're wrong. The headline "
        "result: Expected Calibration Error is significant in every model, Fisher-combined "
        "at 9.3 sigma across 4 Gemini models. Claude at ECE=0.25 provides cross-family validation.\n\n"
        "The direction is consistent: all overconfident. High confidence on wrong answers.\n\n"
        "There's a scaling effect: Pro (ECE=0.23) is significantly better calibrated than the "
        "Flash models (ECE=0.41-0.42). Larger models have confidence surfaces that more closely "
        "track accuracy, but all remain significantly miscalibrated.\n\n"
        "Why this matters for the other results: a model that can't track its own accuracy "
        "can't distinguish perturbations that change moral content from those that change only "
        "surface presentation. Miscalibration is the mechanism that enables the displacement "
        "vulnerabilities measured in the other four tracks."
    )
)
add_body_text(s,
    'Model              ECE      Sigma\n'
    '────────────────────────────────────\n'
    'Flash 2.0          0.414    5.8σ\n'
    'Flash 2.5          0.415    7.0σ\n'
    'Flash 3            0.333    4.5σ\n'
    'Pro 2.5            0.230    2.5σ\n'
    'Claude (cross-fam) 0.250     —\n'
    '────────────────────────────────────\n'
    'Fisher combined            9.3σ',
    top=2.8, size=18)

add_body_text(s,
    'Calibration improves with model scale, but all remain overconfident.',
    top=6.0, size=16, color=AMBER)


# --- Slide 9: Methodology ---
s = add_slide(
    "Methodology: What Makes This Rigorous",
    "Empirical stochastic controls, three-tier data, separation of concerns",
    notes_text=(
        "Three methodological innovations that matter:\n\n"
        "1. EMPIRICAL STOCHASTIC CONTROLS. Every perturbation test includes 3-5 control "
        "replications where the model re-judges identical text. We measure the model's natural "
        "noise, and only displacement beyond that noise counts. This is critical — in preliminary "
        "analyses, apparent >6 sigma invariance violations vanished entirely under empirical "
        "controls. They were measurement artifacts of testing against null=0.\n\n"
        "2. SEPARATION OF CONCERNS. A fixed transformer model generates all perturbation stimuli. "
        "The models under test only judge pre-generated text. This eliminates the self-confirming "
        "loop where a model generates and then evaluates its own test cases.\n\n"
        "3. THREE-TIER DATA. Gold tier: hand-written scenarios with audited ground truth. "
        "Probe tier: synthetic minimal pairs with unambiguous expected behavior. Generated tier: "
        "larger samples from AITA and Dear Abby for statistical power.\n\n"
        "5 models, 2 architecture families (Gemini + Claude). ~8,000 API calls total. "
        "Fisher's method for combining significance across models. Wilson confidence intervals."
    )
)
add_body_text(s,
    '1. Empirical Stochastic Controls\n'
    '   Test against the model\'s own noise, not null=0\n'
    '   (>6σ effects vanished under proper controls)\n\n'
    '2. Separation of Concerns\n'
    '   Fixed transformer generates stimuli; test models only judge\n\n'
    '3. Three-Tier Data Architecture\n'
    '   Gold (hand-audited) → Probe (synthetic) → Generated (AITA/Dear Abby)\n\n'
    '5 models × 2 families × 5 tracks ≈ 8,000 API calls',
    top=2.8, size=18)


# --- Slide 10: Implications ---
s = add_slide(
    "Implications for Alignment",
    "",
    notes_text=(
        "Three takeaways for alignment:\n\n"
        "1. INVARIANCE FAILURES ARE SELECTIVE, NOT TOTAL. Models have learned some symmetries "
        "(gender, order) but not others (framing, emotion). This means the problem is specific "
        "and diagnosable, not a global inability. That's actually encouraging — it means targeted "
        "interventions are possible.\n\n"
        "2. SYCOPHANCY IS ADDRESSABLE. Claude's zero wrong-flip rate proves that training "
        "methodology can eliminate sycophancy. But Claude's simultaneous vulnerability to "
        "emotional anchoring (worst recovery at 20%) shows that fixing one perturbation type "
        "doesn't guarantee robustness to others. Comprehensive evaluation across multiple "
        "perturbation types is necessary.\n\n"
        "3. PROMPT ENGINEERING HAS A CEILING. The ~38% recovery rate sets a practical bound. "
        "For the remaining 62%, architectural interventions — enforcing gauge invariance at the "
        "representation level, not just the prompt level — are likely necessary.\n\n"
        "The geometric framework transforms the alignment question from 'is this model safe?' "
        "to 'which specific invariances does this model maintain, and which does it violate?' "
        "That's a structured diagnostic that guides intervention."
    )
)
add_body_text(s,
    '1. Invariance failures are selective, not total\n'
    '   → Targeted interventions are possible\n\n'
    '2. Sycophancy is addressable (Claude: 0% wrong flip)\n'
    '   → But fixing one vulnerability doesn\'t fix all\n\n'
    '3. Prompt engineering has a ceiling (~38% recovery)\n'
    '   → Architectural interventions needed for the rest\n\n'
    'The question is not "how robust is this model?"\n'
    'but "which invariances does it preserve,\n'
    'and which does it violate?"',
    top=2.6, size=22)


# --- Slide 11: Geometric Reasoning Preview ---
s = add_slide(
    "Looking Ahead: Reasoning as A* on a Manifold",
    "From moral judgment geometry to a general theory of geometric reasoning",
    notes_text=(
        "Where this work leads. The geometric evaluation framework applied to moral judgment "
        "is a special case of a broader thesis: REASONING ITSELF is informed search on a "
        "structured possibility space.\n\n"
        "Think about what A* search does: it navigates a graph using a heuristic that estimates "
        "distance to goal. Now replace the graph with a manifold — a space with curvature, "
        "distance, symmetries, and boundaries. The heuristic becomes a scalar field on that "
        "manifold. Good reasoning follows geodesics; bad reasoning wanders.\n\n"
        "Every failure mode we measured maps to a geometric pathology:\n"
        "- Framing effects = heuristic field warped by salience manipulation\n"
        "- Sycophancy = objective function shifted from truth to approval\n"
        "- Overconfidence = collapsed confidence surface (no gradient signal)\n"
        "- Recovery ceiling = limited metacognitive control of search\n\n"
        "This is the foundation for a new book: 'Geometric Reasoning: From Search to Manifolds.' "
        "The mathematical toolkit is in 'Geometric Methods in Computational Modeling' (Bond, 2026). "
        "The empirical evidence is in the five benchmarks presented today. The theory connects them."
    )
)
add_body_text(s,
    'Reasoning = informed search on a structured space\n\n'
    'Good reasoning   →  follows geodesics\n'
    'Bad reasoning    →  warped heuristic, wrong objective,\n'
    '                      unstable geometry\n\n'
    'Every failure mode we measured is a geometric pathology:\n'
    '  Framing      →  heuristic corruption\n'
    '  Sycophancy   →  proxy-goal capture\n'
    '  Overconfidence → collapsed confidence surface\n'
    '  Recovery ceiling → metacognitive control limit',
    top=2.8, size=20)

add_body_text(s,
    'Bond (2026). Geometric Methods in Computational Modeling. SJSU.',
    top=6.5, size=14, color=TEXT_DIM)


# --- Slide 12: Thank You ---
s = add_slide(
    "Thank You",
    "",
    notes_text=(
        "Thank you. Happy to take questions.\n\n"
        "Key references:\n"
        "- Bond (2026). Geometric Methods in Computational Modeling. SJSU.\n"
        "- Thiele (2026). Does LaBSE encode moral geometry? UCLA.\n"
        "- All benchmark code is available via Kaggle Benchmarks.\n\n"
        "Contact: agi.hpc@gmail.com"
    )
)
add_body_text(s,
    'Andrew H. Bond\n'
    'San Jose State University\n'
    'Department of Computer Engineering\n'
    'agi.hpc@gmail.com\n\n'
    'Lucas Thiele\n'
    'UCLA Cognitive Science\n\n'
    'Benchmark code available on Kaggle Benchmarks\n'
    'Paper: "Selective Invariance Violations in LLM Moral Judgment"',
    top=2.5, size=20, color=TEXT_GRAY)


# ═══════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════

out = Path(__file__).parent / "Measuring_AGI_Presentation.pptx"
prs.save(str(out))
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
